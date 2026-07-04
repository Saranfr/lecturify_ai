"""
Lecturify AI - Document Parser
Extracts text from PDF, DOCX, PPT/PPTX, and TXT files.
Uses pure Python fallbacks (no external binaries) where possible.
"""
import os
import re
import zipfile

try:
    from backend.services.text_preprocessing import is_readable_text as _is_readable_text
except ImportError:
    from services.text_preprocessing import is_readable_text as _is_readable_text
import xml.etree.ElementTree as ET
from typing import List, Optional


def extract_pdf_text(pdf_path: str) -> str:
    """Extract text from PDF using PyMuPDF (fitz) or pdfminer fallback."""
    if not os.path.isfile(pdf_path):
        return ""

    # Try PyMuPDF first (faster)
    try:
        import fitz  # PyMuPDF
        doc = fitz.open(pdf_path)
        text = []
        for page in doc:
            text.append(page.get_text())
        doc.close()
        return "\n".join(text).strip()
    except ImportError:
        pass

    # Fallback: pdfminer
    try:
        from pdfminer.high_level import extract_text as pdfminer_extract
        return pdfminer_extract(pdf_path).strip()
    except ImportError:
        return "[Install PyMuPDF: pip install PyMuPDF]"


def extract_docx_text(docx_path: str) -> str:
    """Extract text from DOCX using python-docx."""
    try:
        from docx import Document
        doc = Document(docx_path)
        return "\n".join(p.text for p in doc.paragraphs if p.text).strip()
    except ImportError:
        return "[Install python-docx: pip install python-docx]"
    except Exception as e:
        return f"[Error parsing DOCX: {e}]"


def _extract_shape_text(shape) -> str:
    """Extract text from a shape (including tables). Filters garbled/symbol font text."""
    parts = []
    def add_part(text: str) -> None:
        t = text.strip()
        if t and _is_readable_text(t):
            parts.append(t)
    if hasattr(shape, "has_table") and shape.has_table:
        try:
            for row in shape.table.rows:
                for cell in row.cells:
                    if cell.text and cell.text.strip():
                        add_part(cell.text)
        except Exception:
            pass
    if hasattr(shape, "text") and shape.text and shape.text.strip():
        add_part(shape.text)
    elif hasattr(shape, "has_text_frame") and shape.has_text_frame:
        try:
            for para in shape.text_frame.paragraphs:
                if para.text and para.text.strip():
                    add_part(para.text)
        except Exception:
            pass
    return " ".join(parts) if parts else ""


def _harvest_strings(data: bytes, min_len: int = 4) -> List[str]:
    """Extract readable string sequences from binary data (UTF-16 and ASCII)."""
    found = []
    # UTF-16-LE (common in Office files)
    try:
        s = data.decode("utf-16-le", errors="ignore")
        for m in re.finditer(r"[\w\u00C0-\u024F\u0400-\u04FF]{3,}", s):
            word = m.group().strip()
            if len(word) >= min_len and not word.isdigit():
                found.append(word)
    except Exception:
        pass
    # ASCII/Latin-1
    try:
        s = data.decode("latin-1", errors="ignore")
        for m in re.finditer(r"[a-zA-Z\u00C0-\u024F]{3,}", s):
            word = m.group().strip()
            if len(word) >= min_len and not word.isdigit():
                found.append(word)
    except Exception:
        pass
    return found


def _extract_ppt_text_raw(ppt_path: str) -> str:
    """
    Best-effort text extraction from .ppt (PowerPoint 97-2003) using OLE + string harvesting.
    Pure Python - no LibreOffice required. Works on any system.
    """
    try:
        import olefile
    except ImportError:
        return "[Install olefile: pip install olefile. Then retry, or save the file as PPTX.]"

    if not olefile.isOleFile(ppt_path):
        return ""
    text_parts = []
    seen = set()
    try:
        ole = olefile.OleFileIO(ppt_path)
        for entry in ole.listdir():
            try:
                data = ole.openstream(entry).read()
                if len(data) < 10:
                    continue
                for w in _harvest_strings(data):
                    w_lower = w.lower()
                    if w_lower not in seen and len(w) > 2:
                        seen.add(w_lower)
                        text_parts.append(w)
            except Exception:
                continue
        ole.close()
    except Exception as e:
        return f"[Error reading PPT: {e}]"
    if not text_parts:
        return "[Could not extract text from .ppt. Save as PPTX (PowerPoint: File > Save As > PPTX) and re-upload.]"
    return " ".join(text_parts[:2000])  # Cap to avoid huge output


def _extract_pptx_text_pure(pptx_path: str) -> str:
    """
    Pure Python PPTX extraction using zipfile + XML. No python-pptx required.
    PPTX is a ZIP of XML files; text is in drawingml <a:t> elements.
    Filters out symbol font / garbled Unicode.
    """
    text_parts = []
    seen = set()
    try:
        with zipfile.ZipFile(pptx_path, "r") as z:
            names = [n for n in z.namelist() if n.endswith(".xml") and ("slide" in n or "notesSlide" in n or "handoutMaster" in n)]
            for name in sorted(names):
                try:
                    data = z.read(name)
                    root = ET.fromstring(data)
                    for elem in root.iter():
                        if elem.tag.endswith("}t") and elem.text:
                            t = elem.text.strip()
                            if t and len(t) >= 2 and _is_readable_text(t) and t.lower() not in seen:
                                seen.add(t.lower())
                                text_parts.append(t)
                except ET.ParseError:
                    pass
                except Exception:
                    pass
    except zipfile.BadZipFile:
        return ""
    except Exception:
        return ""
    return "\n".join(text_parts).strip() if text_parts else ""


def extract_pptx_text(pptx_path: str) -> str:
    """Extract text from PowerPoint (PPT/PPTX). Uses python-pptx if available, else pure Python fallback."""
    if not os.path.isfile(pptx_path):
        return ""
    ext = os.path.splitext(pptx_path)[1].lower()
    if ext == ".ppt":
        result = _extract_pptx_text_pure(pptx_path)
        if result:
            return result
        return _extract_ppt_text_raw(pptx_path)
    result = ""
    try:
        from pptx import Presentation
        prs = Presentation(pptx_path)
        text_parts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                part = _extract_shape_text(shape)
                if part:
                    text_parts.append(part)
            if hasattr(slide, "notes_slide") and slide.notes_slide:
                try:
                    notes_text = slide.notes_slide.notes_text_frame.text
                    if notes_text and notes_text.strip() and _is_readable_text(notes_text.strip()):
                        text_parts.append(notes_text.strip())
                except Exception:
                    pass
        result = "\n".join(text_parts).strip()
    except ImportError:
        result = _extract_pptx_text_pure(pptx_path)
    except Exception:
        result = _extract_pptx_text_pure(pptx_path)
    if not result:
        result = _extract_pptx_text_pure(pptx_path)
    return result


def extract_txt(text_path: str) -> str:
    """Extract text from plain TXT file."""
    if not os.path.isfile(text_path):
        return ""
    try:
        with open(text_path, "r", encoding="utf-8", errors="ignore") as f:
            return f.read().strip()
    except Exception:
        return ""


def extract_document_text(file_path: str) -> str:
    """Route to appropriate parser based on extension."""
    ext = os.path.splitext(file_path)[1].lower()
    if ext == ".pdf":
        return extract_pdf_text(file_path)
    if ext in (".docx", ".doc"):
        return extract_docx_text(file_path)
    if ext in (".pptx", ".ppt"):
        return extract_pptx_text(file_path)
    if ext == ".txt":
        return extract_txt(file_path)
    return ""
